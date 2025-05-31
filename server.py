from mcp.server.fastmcp import FastMCP, Context, Image
from mcp.server.fastmcp.prompts import base
import io
import tempfile
import os
from typing import List, Dict, Optional, Any, Tuple
import base64
import sys
import tempfile
import uuid
import subprocess
import os
from typing import Any, Dict
from pptx.util import Inches, Pt
import oss2
# Create an MCP server for PowerPoint creation with explicit dependencies
mcp = FastMCP(
    "PowerPoint Creator",
    dependencies=["python-pptx","requests"]  # Remove cairosvg dependency
)

# Import pptx after declaring dependencies
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    # Print detailed error to stderr for debugging
    print("Failed to import pptx modules. Make sure python-pptx is installed.", file=sys.stderr)
    raise

# Store presentations in memory during a session
presentations = {}

@mcp.prompt(name="svggenerator", description="""
Generate an SVG image from a natural‑language description, using a simple chain‑of‑thought:
1) Break down your approach in numbered "Thought 1: …", "Thought 2: …" steps.
2) Once you've planned it, output *only* the final SVG markup, wrapped in <svg>…</svg> tags.
""")
def svggenerator_prompt(description: str) -> list[base.Message]:
    return [
        base.UserMessage(
            "You are an expert SVG author.  First think step by step about how to build the image, "
            "then emit just the final SVG markup."
        ),
        base.UserMessage(f"I need an SVG for: {description}")
    ]

@mcp.tool()
def generate_svg(prs_id: str, svg_markup: str, title: str = None, width: float = 6.0) -> str:
    """
    Takes final <svg>…</svg> text, writes to /tmp/, rasterizes via rsvg-convert,
    and inserts as a picture into a new slide.
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"

    # write SVG
    tmp = tempfile.gettempdir()
    fn_base = f"{prs_id}_{uuid.uuid4().hex}"
    svg_path = os.path.join(tmp, fn_base + ".svg")
    with open(svg_path, "w", encoding="utf-8") as f:
        f.write(svg_markup)

    # rasterize using rsvg-convert instead of cairosvg
    png_path = os.path.join(tmp, fn_base + ".png")
    try:
        result = subprocess.run(['rsvg-convert', '-o', png_path, svg_path], 
                                check=True, 
                                capture_output=True, 
                                text=True)
        if result.returncode != 0:
            return f"SVG→PNG error: {result.stderr}"
    except subprocess.SubprocessError as e:
        return f"SVG→PNG error: {str(e)}"
    except FileNotFoundError:
        return "Error: rsvg-convert not found. Please install librsvg with 'brew install librsvg'"

    # insert slide
    prs = presentations[prs_id]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if title:
        tx = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        p  = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True

    left = Inches((10 - width) / 2)
    slide.shapes.add_picture(png_path, left, Inches(2), width=Inches(width))

    # cleanup
    for p in (svg_path, png_path):
        try: os.remove(p)
        except: pass

    return f"Inserted SVG slide at position {len(prs.slides)}"


@mcp.tool()
def create_presentation(title: str) -> str:
    """
    Create a new empty PowerPoint presentation with a title.
    
    Args:
        title: The title/name for the presentation
        
    Returns:
        Presentation ID to use in subsequent operations
    """
    prs = Presentation()
    prs_id = title.replace(" ", "_").lower()
    presentations[prs_id] = prs
    return f"Created presentation: {prs_id}"

@mcp.tool()
def add_title_slide(prs_id: str, title: str, subtitle: Optional[str] = None) -> str:
    """
    Add a title slide to the presentation.
    
    Args:
        prs_id: The presentation ID
        title: The title text
        subtitle: Optional subtitle text
        
    Returns:
        Confirmation message with slide number
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Set subtitle if provided
    if subtitle and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        if subtitle_shape.has_text_frame:
            subtitle_shape.text = subtitle
    
    return f"Added title slide at position {len(prs.slides)}"

@mcp.tool()
def add_content_slide(prs_id: str, title: str, content: List[str]) -> str:
    """
    Add a content slide with bullet points.
    
    Args:
        prs_id: The presentation ID
        title: The slide title
        content: List of bullet points
        
    Returns:
        Confirmation message with slide number
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if hasattr(slide.shapes, 'title') and slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title
    
    # Find content placeholder
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # Body placeholder
            content_placeholder = shape
            break
    
    # If no content placeholder found, add a textbox
    if not content_placeholder:
        content_placeholder = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
    
    # Add bullet points
    tf = content_placeholder.text_frame
    tf.clear()  # Clear any existing text
    
    # Add first bullet point
    if content and len(content) > 0:
        p = tf.paragraphs[0]
        p.text = content[0]
        p.level = 0
        
        # Add remaining bullet points
        for bullet_text in content[1:]:
            p = tf.add_paragraph()
            p.text = bullet_text
            p.level = 0
    
    return f"Added content slide at position {len(prs.slides)}"

@mcp.tool()
def add_section_slide(prs_id: str, section_title: str, background_color: Optional[str] = None) -> str:
    """
    Add a section divider slide with a large title.
    
    Args:
        prs_id: The presentation ID
        section_title: The section title text
        background_color: Optional hex color code (e.g. "#FF0000" for red)
        
    Returns:
        Confirmation message with slide number
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color if provided
    if background_color and background_color.startswith("#"):
        try:
            bg_color = background_color.lstrip("#")
            r, g, b = tuple(int(bg_color[i:i+2], 16) for i in (0, 2, 4))
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"Error setting background color: {e}", file=sys.stderr)
    
    # Add large centered title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = section_title
    
    # Format the title text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True

    if background_color and background_color.lower() in ["#000000", "#0000ff", "#800080", "#000080"]:
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    return f"Added section slide at position {len(prs.slides)}"

@mcp.tool()
def add_image_slide(prs_id: str, title: str, image_path: str, image_description: str) -> str:
    """
    Add a slide with an image and title.
    
    Args:
        prs_id: The presentation ID
        title: The slide title
        image_path: Path to the image file or URL
        image_description: Description of the image (alt text)
        
    Returns:
        Confirmation message with slide number
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add image
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    
    try:
        if image_path.startswith(('http://', 'https://')):
            import requests
            import tempfile
            
            response = requests.get(image_path, stream=True)
            response.raise_for_status()  
            
            img_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
            img_temp_path = img_temp.name
            
            with open(img_temp_path, 'wb') as f:
                for chunk in response.iter_content(1024):
                    f.write(chunk)
            
            pic = slide.shapes.add_picture(img_temp_path, left, top, width=width)
            
            try:
                os.unlink(img_temp_path)
            except:
                pass
        else:
            if not os.path.exists(image_path):
                return f"Error: Image file '{image_path}' not found"
            pic = slide.shapes.add_picture(image_path, left, top, width=width)
        
        try:
            if hasattr(pic, 'shape_properties'):
                pic.shape_properties.has_text_frame = True
                pic.text_frame.text = image_description
        except:
            alt_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(0.1), Inches(0.1))
            alt_box.text_frame.text = f"Alt text: {image_description}"
    except Exception as e:
        print(f"Error adding image: {e}", file=sys.stderr)
        return f"Error adding image: {str(e)}"
    
    return f"Added image slide at position {len(prs.slides)}"

@mcp.tool()
def add_table_slide(prs_id: str, title: str, headers: List[str], rows: List[List[str]]) -> str:
    """
    Add a slide with a table.
    
    Args:
        prs_id: The presentation ID
        title: The slide title
        headers: List of column headers
        rows: List of rows, where each row is a list of cell values
        
    Returns:
        Confirmation message with slide number
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Calculate table dimensions
    num_rows = len(rows) + 1  # +1 for header row
    num_cols = len(headers)
    
    # Create table
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.5 * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    
    # Set column widths
    col_width = width / num_cols
    for i in range(num_cols):
        table.columns[i].width = col_width
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add data rows
    for i, row in enumerate(rows):
        for j, cell_value in enumerate(row):
            if j < num_cols:  
                cell = table.cell(i + 1, j)
                cell.text = cell_value
    
    return f"Added table slide at position {len(prs.slides)}"

@mcp.tool()
def save_presentation(prs_id: str, output_path: str) -> str:
    """
    Save the presentation to a file.
    
    Args:
        prs_id: The presentation ID
        output_path: File path where to save the presentation
        
    Returns:
        Confirmation message with file path
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    
    try:
        # 生成临时文件名（使用UUID避免冲突）
        file_name = f"{prs_id}_{uuid.uuid4().hex[:8]}.pptx"
        tmp_path = os.path.join("/tmp", file_name)
        
        # 保存到临时文件
        prs.save(tmp_path)
        file_size = os.path.getsize(tmp_path)
        
        # 获取OSS配置（使用环境变量）
        endpoint = os.getenv("OSS_ENDPOINT")
        bucket_name = os.getenv("OSS_BUCKET_NAME")
        access_key = os.getenv("OSS_ACCESS_KEY")
        secret_key = os.getenv("OSS_SECRET_KEY")
        
        if not all([endpoint, bucket_name, access_key, secret_key]):
            return "OSS configuration missing"
        
        # 创建OSS客户端
        auth = oss2.Auth(access_key, secret_key)
        bucket = oss2.Bucket(auth, endpoint, bucket_name)
        
        # 上传文件
        oss_object_name = f"presentations/{file_name}"
        bucket.put_object_from_file(oss_object_name, tmp_path)
        
        # 生成可下载URL（1小时有效）
        download_url = bucket.sign_url('GET', oss_object_name, 3600)
        
        # 清理临时文件
        os.remove(tmp_path)
        
        return f"PPT已保存！下载地址：{download_url}"
    
    except Exception as e:
        error_msg = f"Error saving presentation: {str(e)}"
        print(error_msg, file=sys.stderr)
        return error_msg

@mcp.tool()
def get_presentation_download_link(prs_id: str) -> str:
    """
    Get a download link or base64 representation for the presentation.
    
    Args:
        prs_id: The presentation ID
        
    Returns:
        Base64 encoded presentation data that can be downloaded
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    try:
        tmp_file = os.path.join("/tmp", f"{prs_id}.pptx")
        presentations[prs_id].save(tmp_file)
        
        with open(tmp_file, "rb") as f:
            file_data = f.read()
        
        base64_data = base64.b64encode(file_data).decode("utf-8")
        
        download_uri = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{base64_data}"
        
        return f"""
Presentation ready for download.
File size: {len(file_data) / 1024:.2f} KB

To download:
1. Copy the base64 string below
2. Use a base64-to-file converter or
3. In a browser console, use: 
   const a = document.createElement('a');
   a.href = '{download_uri}';
   a.download = '{prs_id}.pptx';
   a.click();

Base64 data: {base64_data[:100]}...
"""
    except Exception as e:
        return f"Error creating download link: {str(e)}"
    
@mcp.tool()
def get_presentation_info(prs_id: str) -> str:
    """
    Get information about the presentation.
    
    Args:
        prs_id: The presentation ID
        
    Returns:
        Information about the presentation
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    slide_count = len(prs.slides)
    
    info = f"Presentation: {prs_id}\n"
    info += f"Number of slides: {slide_count}\n"
    info += "Slide layouts available:\n"
    
    for i, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
        info += f"  - {layout_name}\n"
    
    return info

@mcp.resource("presentation://{prs_id}/outline")
def get_presentation_outline(prs_id: str) -> str:
    """
    Get a text representation of the presentation structure.
    
    Args:
        prs_id: The presentation ID
        
    Returns:
        Text outline of the presentation
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    
    outline = f"# Presentation: {prs_id}\n\n"
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_type = "Unknown"
        slide_title = "Untitled"
        
        # Try to determine slide type and title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide_title = slide.shapes.title.text
        
        # Guess slide type based on layout
        if hasattr(slide.slide_layout, 'name') and slide.slide_layout.name:
            slide_type = slide.slide_layout.name
        elif i == 0 and len(prs.slides) > 0:
            slide_type = "Title Slide"
        
        outline += f"## Slide {slide_num}: {slide_title}\n"
        outline += f"Type: {slide_type}\n"
        
        # Add content summary
        content_summary = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape != getattr(slide.shapes, 'title', None):
                text_summary = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                content_summary.append(f"- Text: {text_summary}")
            elif hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE.PICTURE:
                content_summary.append("- Image")
            elif hasattr(shape, 'shape_type') and hasattr(shape, 'table') and shape.shape_type == MSO_SHAPE.TABLE:
                content_summary.append(f"- Table ({shape.table.rows_count}x{shape.table.columns_count})")
        
        if content_summary:
            outline += "Content:\n" + "\n".join(content_summary) + "\n"
        
        outline += "\n"
    
    return outline

@mcp.tool()
def remove_slide(prs_id: str, slide_index: int) -> str:
    """
    Remove a slide from the presentation.
    
    Args:
        prs_id: The presentation ID
        slide_index: The index of the slide to remove (1-based)
        
    Returns:
        Confirmation message
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    
    # Convert 1-based index to 0-based index
    zero_based_index = slide_index - 1
    
    if zero_based_index < 0 or zero_based_index >= len(prs.slides):
        return f"Error: Slide index {slide_index} is out of range (valid range: 1-{len(prs.slides)})"
    
    try:
        # This is a simplified approach that might not work in all python-pptx versions
        # Get XML ids needed for removal
        if hasattr(prs.slides, '_sldIdLst'):
            slide_id = prs.slides._sldIdLst[zero_based_index].attrib['id']
            slide_rel = prs.slides._sldIdLst[zero_based_index].attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
            
            # Remove the slide
            prs.part.drop_rel(slide_rel)
            del prs.slides._sldIdLst[zero_based_index]
        else:
            return f"Cannot remove slide - this functionality requires direct XML manipulation that may not be supported in this version of python-pptx"
    except Exception as e:
        print(f"Error removing slide: {e}", file=sys.stderr)
        return f"Error removing slide: {str(e)}"
    
    return f"Removed slide at position {slide_index}"

@mcp.tool()
def export_to_base64(prs_id: str) -> str:
    """
    Export the presentation as a base64-encoded string.
    
    Args:
        prs_id: The presentation ID
        
    Returns:
        Base64-encoded presentation (first 100 characters shown)
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    prs = presentations[prs_id]
    
    try:
        # Save to a bytes buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        # Encode to base64
        base64_data = base64.b64encode(buffer.read()).decode('utf-8')
        
        # Return a shortened version to confirm it worked
        return f"Base64 encoded presentation (showing first 100 chars): {base64_data[:100]}..."
    except Exception as e:
        error_msg = f"Error exporting presentation to base64: {str(e)}"
        print(error_msg, file=sys.stderr)
        return error_msg

@mcp.tool()
def create_presentation_with_outline(title: str, outline: Dict[str, Any]) -> str:
    """
    Create a complete presentation from a structured outline in one operation.
    
    Args:
        title: Presentation title
        outline: Structured outline with slides data
        
    Example outline:
    {
        "title_slide": {"title": "My Presentation", "subtitle": "Subtitle"},
        "slides": [
            {"type": "content", "title": "Slide 1", "content": ["Point 1", "Point 2"]},
            {"type": "section", "title": "Section Break", "color": "#0066CC"},
            {"type": "svg", "title": "Chart", "description": "Bar chart showing sales data"},
            {"type": "table", "title": "Data", "headers": ["A", "B"], "rows": [["1", "2"]]}
        ]
    }
    """
    prs_id = title.replace(" ", "_").lower()
    
    # Create presentation
    prs = Presentation()
    presentations[prs_id] = prs
    
    # Add title slide if specified
    if "title_slide" in outline:
        ts = outline["title_slide"]
        add_title_slide(prs_id, ts["title"], ts.get("subtitle"))
    
    # Process all slides in batch
    for slide_data in outline.get("slides", []):
        slide_type = slide_data["type"]
        
        if slide_type == "content":
            add_content_slide(prs_id, slide_data["title"], slide_data["content"])
        elif slide_type == "section":
            add_section_slide(prs_id, slide_data["title"], slide_data.get("color"))
        elif slide_type == "svg":
            # Generate SVG internally without exposing steps
            _generate_svg_internal(prs_id, slide_data["description"], slide_data["title"])
        elif slide_type == "table":
            add_table_slide(prs_id, slide_data["title"], slide_data["headers"], slide_data["rows"])
    
    return f"Created presentation '{prs_id}' with {len(outline.get('slides', []))} slides"

def _generate_svg_internal(prs_id: str, description: str, title: str) -> None:
    """Internal SVG generation - no LLM chain-of-thought, just direct generation"""
    # This would internally call your SVG generation logic
    # without exposing the step-by-step process to the LLM
    pass

@mcp.tool()  
def add_multiple_content_slides(prs_id: str, slides_data: List[Dict[str, Any]]) -> str:
    """
    Add multiple content slides in one operation.
    
    Args:
        prs_id: Presentation ID
        slides_data: List of slide data dictionaries
        
    Example:
    [
        {"title": "Introduction", "content": ["Welcome", "Overview", "Goals"]},
        {"title": "Main Points", "content": ["Point A", "Point B", "Point C"]},
        {"title": "Conclusion", "content": ["Summary", "Next Steps"]}
    ]
    """
    if prs_id not in presentations:
        return f"Error: Presentation '{prs_id}' not found"
    
    added_count = 0
    for slide_data in slides_data:
        result = add_content_slide(prs_id, slide_data["title"], slide_data["content"])
        if not result.startswith("Error"):
            added_count += 1
    
    return f"Added {added_count} content slides to presentation"

@mcp.tool()
def create_quick_presentation(title: str, slide_titles: List[str], content_per_slide: List[List[str]]) -> str:
    """
    Quickly create a standard content presentation with multiple slides.
    
    Args:
        title: Presentation title
        slide_titles: List of slide titles
        content_per_slide: List of bullet points for each slide
    """
    prs_id = title.replace(" ", "_").lower()
    
    # Create and build presentation in one go
    prs = Presentation()
    presentations[prs_id] = prs
    
    # Title slide
    add_title_slide(prs_id, title, "Generated Presentation")
    
    # Content slides
    for i, slide_title in enumerate(slide_titles):
        content = content_per_slide[i] if i < len(content_per_slide) else []
        add_content_slide(prs_id, slide_title, content)
    
    return f"Created presentation '{prs_id}' with {len(slide_titles)} content slides"

def main() -> None:
    """CLI wrapper so we can run `python -m mcp_server_ppt_maker`"""
    mcp.run()
